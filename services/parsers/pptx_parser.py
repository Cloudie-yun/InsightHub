from pathlib import Path

from services.parsers.utils import clean_extracted_text


def _dependency_error(package_name):
    return {
        "code": "missing_dependency",
        "message": f"Required parser dependency '{package_name}' is not installed.",
        "details": {"package": package_name},
    }


def _segment_sort_key(segment):
    return (
        segment.get("source_index", 0),
        segment.get("block_index", 0),
        segment.get("paragraph_index", 0),
        segment.get("segment_id", ""),
    )


def parse_pptx(file_path):
    try:
        from pptx import Presentation
    except Exception:
        return {
            "segments": [],
            "metadata": {},
            "errors": [_dependency_error("python-pptx")],
        }

    file_path = Path(file_path)
    segments = []
    metadata = {
        "source_path": str(file_path),
        "slide_count": 0,
    }

    try:
        presentation = Presentation(file_path)
        metadata["slide_count"] = len(presentation.slides)

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_text_parts = []
            for shape in sorted(slide.shapes, key=lambda current_shape: current_shape.shape_id):
                if not hasattr(shape, "text"):
                    continue
                text = clean_extracted_text(shape.text or "")
                if text:
                    slide_text_parts.append(text)

            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = clean_extracted_text(slide.notes_slide.notes_text_frame.text or "")
                if notes_text:
                    slide_text_parts.append(notes_text)

            if not slide_text_parts:
                continue

            slide_text = clean_extracted_text("\n\n".join(slide_text_parts))
            segments.append(
                {
                    "segment_id": f"pptx-slide-{slide_index}",
                    "source_type": "slide",
                    "source_index": slide_index,
                    "text": slide_text,
                    "metadata": {
                        "shape_text_count": len(slide_text_parts),
                        "char_count": len(slide_text),
                        "parser_adapter": "pptx",
                    },
                }
            )

        segments.sort(key=_segment_sort_key)
        return {"segments": segments, "metadata": metadata, "errors": []}
    except Exception as exc:
        return {
            "segments": [],
            "metadata": metadata,
            "errors": [
                {
                    "code": "pptx_parse_error",
                    "message": "Unable to parse PPTX file.",
                    "details": {"exception": str(exc), "source_path": str(file_path)},
                }
            ],
        }
